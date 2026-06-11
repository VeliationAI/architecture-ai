#!/usr/bin/env python3
"""Generate Architecture AI Studio pitch deck (.pptx)."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand palette (matches studio dark theme)
BG_DARK = RGBColor(0x0C, 0x0C, 0x12)
BG_CARD = RGBColor(0x16, 0x16, 0x22)
ACCENT = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_LIGHT = RGBColor(0xA7, 0x8B, 0xFA)
TEXT_PRIMARY = RGBColor(0xF4, 0xF4, 0xF5)
TEXT_MUTED = RGBColor(0xA1, 0xA1, 0xAA)
SUCCESS = RGBColor(0x34, 0xD3, 0x99)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT = Path(__file__).resolve().parent.parent / "doc" / "Architecture-AI-Studio.pptx"


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=Inches(0), height=Inches(0.06)) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_glow_orb(slide, left, top, size, alpha=0.15) -> None:
    orb = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    orb.fill.solid()
    orb.fill.fore_color.rgb = ACCENT
    orb.fill.transparency = 1 - alpha
    orb.line.fill.background()


def text_box(slide, left, top, width, height, text, size=18, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullets(slide, left, top, width, height, items, size=16, color=TEXT_MUTED):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.level = 0
    return tf


def feature_card(slide, left, top, w, h, title, desc):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = RGBColor(0x2A, 0x2A, 0x36)
    card.line.width = Pt(1)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), h)
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    text_box(slide, left + Inches(0.25), top + Inches(0.2), w - Inches(0.4), Inches(0.4), title, size=14, bold=True, color=ACCENT_LIGHT)
    text_box(slide, left + Inches(0.25), top + Inches(0.55), w - Inches(0.4), h - Inches(0.65), desc, size=11, color=TEXT_MUTED)


def title_slide(prs, title, subtitle="", badge=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_bar(slide)
    add_glow_orb(slide, Inches(9), Inches(-1), Inches(4), 0.12)
    add_glow_orb(slide, Inches(-1), Inches(5), Inches(3), 0.08)

    if badge:
        text_box(slide, Inches(0.9), Inches(1.2), Inches(6), Inches(0.4), badge, size=12, color=ACCENT_LIGHT)

    text_box(slide, Inches(0.9), Inches(1.8), Inches(11), Inches(1.5), title, size=40, bold=True, color=TEXT_PRIMARY)
    if subtitle:
        text_box(slide, Inches(0.9), Inches(3.4), Inches(10), Inches(1.2), subtitle, size=18, color=TEXT_MUTED)


def section_slide(prs, section_num, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_bar(slide, top=Inches(3.2), height=Inches(0.04))
    text_box(slide, Inches(0.9), Inches(2.4), Inches(2), Inches(0.6), f"0{section_num}", size=48, bold=True, color=ACCENT)
    text_box(slide, Inches(0.9), Inches(3.5), Inches(11), Inches(1), title, size=32, bold=True, color=TEXT_PRIMARY)
    if subtitle:
        text_box(slide, Inches(0.9), Inches(4.5), Inches(10), Inches(0.8), subtitle, size=16, color=TEXT_MUTED)


def content_slide(prs, title, bullets=None, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_bar(slide)
    text_box(slide, Inches(0.9), Inches(0.5), Inches(11), Inches(0.7), title, size=28, bold=True, color=TEXT_PRIMARY)
    if bullets:
        add_bullets(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5), bullets, size=17)
    if footer:
        text_box(slide, Inches(0.9), Inches(6.8), Inches(11), Inches(0.4), footer, size=11, color=TEXT_MUTED)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 — Title
    title_slide(
        prs,
        "Architecture AI Studio",
        "Turn use cases into explainable target architectures\nwith platform-native best practice review.",
        "VELIATION AI  ·  MIT OPEN SOURCE",
    )
    s = prs.slides[-1]
    text_box(s, Inches(0.9), Inches(5.0), Inches(8), Inches(0.5), "architecture-ai.onrender.com", size=15, color=SUCCESS)

    # 2 — Problem
    section_slide(prs, 1, "The problem", "Diagrams don't ship. Decisions get lost.")
    content_slide(
        prs,
        "Why existing tools fall short",
        [
            "●  Architecture lives in slides — not typed, versioned, or exportable",
            "●  One recommendation ignores speed vs governance vs scale tradeoffs",
            "●  Platform features change monthly — review checklists go stale",
            "●  AI chat produces boxes without why, benefit, tradeoff, or risk",
        ],
    )

    # 3 — Solution
    section_slide(prs, 2, "Our solution", "Architecture intelligence, not diagramming.")
    slide = content_slide(prs, "End-to-end workflow", footer="")
    flow_items = [
        ("Intake", "Structured requirements"),
        ("Portfolio", "3 strategic variants"),
        ("Canvas", "Interactive diagram"),
        ("Review", "7D WAF scoring"),
        ("Approve", "Sign-off workflow"),
        ("Export", "JSON · Mermaid · IaC"),
    ]
    x = Inches(0.6)
    for label, sub in flow_items:
        feature_card(slide, x, Inches(1.6), Inches(1.85), Inches(1.3), label, sub)
        x += Inches(2.05)

    # 4 — Features grid
    section_slide(prs, 3, "Key capabilities")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_bar(slide)
    text_box(slide, Inches(0.9), Inches(0.5), Inches(11), Inches(0.7), "What makes it different", size=28, bold=True)
    cards = [
        ("Multi-variant portfolio", "MVP, governed, and scale-ready designs with scores"),
        ("Typed graph schema", "Nodes, edges, zones, rationale — not pixels"),
        ("7 review dimensions", "Security through explainability"),
        ("Platform knowledge", "Versioned Databricks & AWS registries"),
        ("Data models", "Conceptual → star schema + transforms"),
        ("MCP + exports", "Agent tools & engineering artifacts"),
    ]
    positions = [(0.9, 1.4), (4.6, 1.4), (8.3, 1.4), (0.9, 3.5), (4.6, 3.5), (8.3, 3.5)]
    for (l, t), (title, desc) in zip(positions, cards):
        feature_card(slide, Inches(l), Inches(t), Inches(3.4), Inches(1.75), title, desc)

    # 5 — Architecture
    content_slide(
        prs,
        "Technical architecture",
        [
            "apps/web        —  Next.js 15 studio, React Flow canvas, REST APIs",
            "packages/core   —  Graph schema, LLM client, portfolio, review, export",
            "packages/catalog — Platform knowledge, icons, domain packs",
            "packages/mcp    —  MCP server for Cursor / Claude / custom agents",
            "",
            "Deploy: Render (production)  ·  Mock LLM for offline demo",
        ],
    )

    # 6 — Review dimensions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_bar(slide)
    text_box(slide, Inches(0.9), Inches(0.5), Inches(11), Inches(0.7), "Well-architected review engine", size=28, bold=True)
    dims = ["Security", "Reliability", "Performance", "Cost", "Operations", "Governance", "Explainability"]
    for i, dim in enumerate(dims):
        col = i % 4
        row = i // 4
        left = Inches(0.9 + col * 3.1)
        top = Inches(1.6 + row * 1.5)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1.5)
        text_box(slide, left, top + Inches(0.35), Inches(2.8), Inches(0.5), dim, size=16, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)

    # 7 — Workspace UX
    content_slide(
        prs,
        "Studio workspace",
        [
            "Design     →  Canvas · Variants · Compare · Data model",
            "Workflow   →  Improve · Review · Approve",
            "Deliver    →  Summary · Export (JSON, Mermaid, Terraform, dbt, ADF)",
            "",
            "Left nav opens content in the main panel — intuitive, mobile-friendly",
        ],
    )

    # 8 — Live demo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_glow_orb(slide, Inches(4), Inches(1.5), Inches(5), 0.18)
    add_accent_bar(slide)
    text_box(slide, Inches(0.9), Inches(2.2), Inches(11), Inches(0.8), "Try it live", size=36, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, Inches(0.9), Inches(3.2), Inches(11), Inches(0.6), "architecture-ai.onrender.com", size=22, color=SUCCESS, align=PP_ALIGN.CENTER)
    text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(1.2), "No API key required · Mock LLM mode · Free on Render", size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # 9 — Roadmap
    content_slide(
        prs,
        "Roadmap",
        [
            "✓  Multi-variant portfolio, review, approval, export, MCP",
            "✓  Databricks deep knowledge (Genie, Unity Catalog, MLflow)",
            "✓  AWS well-architected basics",
            "→  Azure, GCP, Snowflake rule packs",
            "→  LLM app patterns (routing, grounding, evaluation)",
            "→  Product demo video on landing page",
        ],
    )

    # 10 — Team
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_bar(slide)
    text_box(slide, Inches(0.9), Inches(0.5), Inches(11), Inches(0.7), "Team", size=28, bold=True)
    feature_card(slide, Inches(1.5), Inches(2.0), Inches(4.5), Inches(1.8), "Akhil Vydyula", "Core contributor\nlinkedin.com/in/akhilvydyula")
    feature_card(slide, Inches(7.0), Inches(2.0), Inches(4.5), Inches(1.8), "Sai Sankara Thamma", "Core contributor\nlinkedin.com/in/sankara-reddy-thamma-18a6a6ba")
    text_box(slide, Inches(0.9), Inches(5.2), Inches(11), Inches(0.5), "MIT License  ·  github.com/VeliationAI/architecture-ai", size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # 11 — Thank you
    title_slide(prs, "Thank you", "Architecture intelligence, not diagramming.", "QUESTIONS?")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
